from io import BytesIO

import av
import pptx
from pptx.enum.action import PP_ACTION
from pptx.util import Inches

from render import run_and_rename


def extract_first_frame_to_bytes(video_path):
    """Extract the first frame from a video and return it as bytes."""
    container = av.open(str(video_path))

    # Hole den ersten Frame
    for frame in container.decode(video=0):
        # Konvertiere zu PIL Image
        img = frame.to_image()
        # Speichere in BytesIO
        img_bytes = BytesIO()
        img.save(img_bytes, format='JPEG')
        img_bytes.seek(0)
        container.close()
        return img_bytes

    container.close()
    return None


def generate_pptx(filename, scene_name):
    quality = '2160p60'

    video_files = run_and_rename(filename, scene_name, quality)

    prs = pptx.Presentation()
    prs.slide_width = Inches(16)       # 9-Teil
    prs.slide_height = Inches(9)
    slide_layout = prs.slide_layouts[6]  # Title Only layout


    for video_file in video_files:
        slide = prs.slides.add_slide(slide_layout)
        slide.name = f'{scene_name}_{video_file}'

        poster_bytes = extract_first_frame_to_bytes(video_file)

        left = 0
        top = 0
        width = prs.slide_width
        height = prs.slide_height

        slide.shapes.add_movie(video_file, left, top, width, height,
                               mime_type='video/mp4', poster_frame_image=poster_bytes)

    prs.save('presentation.pptx')

    print('Presentation saved as presentation.pptx')
    print('What you now need to do:')
    print('1. Import the presentation.pptx into Keynote.')
    print('2. For each slide, set the video to "Start on Click".')
    print('3. Set each Slide to "Automatically" advance.')


if __name__ == '__main__':
    scene_name = 'Task641036'
    filename = 'math_olympiade.py'
    generate_pptx(filename, scene_name)
